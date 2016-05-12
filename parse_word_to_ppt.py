#!/usr/bin/python
# -*- coding: UTF-8 -*-

import traceback
import os, sys

# check os
import platform

if platform.system() == "Windows":
    sys.path.append(os.path.join('python_package', 'lib_for_win'))

sys.path.append(os.path.join('python_package'))


#for date
import datetime
today = str(datetime.date.today())

# for thread
import thread
import time

# for log
from  log_config import *

# for pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches,Cm,Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.dml import MSO_COLOR_TYPE

EMPTY_SLIDE=6
DEFAULT_SLIDE=0

# for word
from docx import Document
from docx.shared import Inches as doc_Inches
from docx.shared import Pt as doc_Pt

document = Document()

TITLE_PREFIX="title:"

ENCODING='utf8'

WIDEST_LINE_IN_PPT=13
LANGEST_LINE_IN_PPT=6

def query_yes_no(question, default="yes"):
    """Ask a yes/no question via raw_input() and return their answer.

    "question" is a string that is presented to the user.
    "default" is the presumed answer if the user just hits <Enter>.
        It must be "yes" (the default), "no" or None (meaning
        an answer is required of the user).

    The "answer" return value is True for "yes" or False for "no".
    """
    valid = {"yes": True, "y": True, "ye": True,
             "no": False, "n": False}
    if default is None:
        prompt = " [y/n] "
    elif default == "yes":
        prompt = " [Y/n] "
    elif default == "no":
        prompt = " [y/N] "
    else:
        raise ValueError("invalid default answer: '%s'" % default)

    while True:
        sys.stdout.write(question + prompt)
        choice = raw_input().lower()
        if default is not None and choice == '':
            return valid[default]
        elif choice in valid:
            return valid[choice]
        else:
            sys.stdout.write("Please respond with 'yes' or 'no' "
                             "(or 'y' or 'n').\n")

class DocProduct:
    def __init__(self):
        self.song_list_ = {}

    def GetSongList(self, song_title_list):
        ret_song_dict = {}
        find_all = True
        for title in song_title_list:
            if title in self.song_list_:
                log.debug("get song:{}, paragraph.num:{}".format(title.encode(ENCODING), len(self.song_list_[title]) ))
                ret_song_dict[title] = self.song_list_[title]
            else:
                log.warning(u"无法找到\t\"{}\"\t的歌词".format(title))
                find_all = False
        return (find_all, ret_song_dict)


    def WriteDoc(self, song_title_list, song_lyric_dict):
        document = Document("template.docx")
        for index, title in enumerate(song_title_list):
            #p = document.add_paragraph(TITLE_PREFIX + title).bold = True
            p = ""
            if index == 0:
                p = document.paragraphs[0]
            else:
                p = document.add_paragraph()

            run = p.add_run(TITLE_PREFIX + title)
            font = run.font
            font.bold = True
            #font.name = 'SimSun'
            font.name = 'Microsoft YaHei'
            font.size = doc_Pt(12)
            #p = document.add_paragraph()

            song_lyric_paragraph_list = song_lyric_dict.get(title)

            if song_lyric_paragraph_list == None:
                continue

            for song_lyric_paragraph in song_lyric_paragraph_list:
                for song_lyric_ in song_lyric_paragraph:
                    p = document.add_paragraph()
                    run = p.add_run(song_lyric_)
                    font = run.font
                    font.size = doc_Pt(12)
                p = document.add_paragraph()

        document.save(today + '.docx')
        return

    def ReadDoc(self, name):
        log.debug ("DocProduct:ReadDoc, name:{} ".format(name.encode(ENCODING)))
        document = Document(name)
        in_the_song = False
        song_tile = ""
        song_lyric_paragraph = []
        song_lyric_paragraph_list = []
        for paragraph in document.paragraphs:
            text = paragraph.text
            if TITLE_PREFIX[:-1] in text.lower():
                log.debug("enter the song:{}".format(text.encode(ENCODING)))
                if song_tile != "":
                    self.song_list_[song_tile] = song_lyric_paragraph_list
                    log.debug("finish parse song:{} with {} paragraph"
                             .format(song_tile.encode(ENCODING), len(self.song_list_[song_tile])))
                    song_lyric_paragraph = []
                    song_lyric_paragraph_list = []
                    song_tile = ""

                song_tile = text[len(TITLE_PREFIX):]
                log.debug ("get song_tile:{}".format(song_tile.encode(ENCODING)))
                in_the_song = True
            elif in_the_song == True:
                if text != "":
                    log.debug("enter the song lyric paragraph:{}".format(text.encode(ENCODING)))
                    song_lyric_paragraph.append(text)
                elif len(song_lyric_paragraph):
                    log.debug("leave the song lyric paragraph:{}".format(text.encode(ENCODING)))
                    song_lyric_paragraph_list.append(song_lyric_paragraph)
                    song_lyric_paragraph = []

        if song_tile != "":
            self.song_list_[song_tile] = song_lyric_paragraph_list
            log.debug("finish parse song:{} with {} paragraph"
                     .format(song_tile.encode(ENCODING), len(song_lyric_paragraph_list)))
            song_lyric_paragraph_list = []
            song_tile = ""





class PPTProduct:
    def __init__(self):
        return

    def AddLyric(self, slide, content_list):
        log.debug ("PPTProduct:AddLyric, content:{} ".format(content_list))

        shapes = slide.shapes

        left = Cm(0.36)
        top = Cm(3.4)

        if 4 < len(content_list) and len(content_list) < LANGEST_LINE_IN_PPT:
            top = Cm(1.72)
        elif len(content_list) == LANGEST_LINE_IN_PPT:
            top = Cm(0.72)

        width = Cm(24.7)
        height = Cm(3.0)
        shape = shapes.add_textbox(left, top, width, height);
        # shape.text = "Hello, World!\n test\n test\n"

        text_frame = shape.text_frame
        text_frame.clear()  # remove any existing paragraphs, leaving one empty one

        # p = text_frame.paragraphs[0]
        # p.text = paragraph_strs[0]

        for para_str in content_list:
            p = text_frame.add_paragraph()
            p.text = para_str
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing=Pt(75)
            font = p.font
            font.name = 'Microsoft YaHei'
            #font.size = Pt(56)
            if (len(para_str) > WIDEST_LINE_IN_PPT) :
                font.size = Pt((24.7/len(para_str)) * 28.3464) # 1cm = 28.3464pt
            else:
                font.size = Pt(56)
            font.bold = False
            font.italic = None  # cause value to be inherited from theme
            font.color.rgb = RGBColor(0xFD, 0xBF, 0x2D)
            # font.color.theme_color = MSO_THEME_COLOR.ACCENT_1


    def AddTitle(self, slide, title):
        log.debug ("PPTProduct:AddTitle, title:{}".format(title.encode(ENCODING)))

        shapes = slide.shapes
        left = Cm(0)
        top = Cm(18)
        width = Cm(25.15)
        height = Cm(1.03)
        shape = shapes.add_textbox(left, top, width, height);

        text_frame = shape.text_frame
        text_frame.clear()  # remove any existing paragraphs, leaving one empty one

        p = text_frame.paragraphs[0]
        p.text = title

        p.alignment = PP_ALIGN.RIGHT
#        p.line_spacing=Pt(75)
        font = p.font
        font.name = 'Microsoft YaHei'
        font.size = Pt(20)
        font.bold = True
        font.italic = None  # cause value to be inherited from theme
        font.color.rgb = RGBColor(0xFD, 0xBF, 0x2D)

    def CreateSlide(self, song_title_list, song_lyric_paragraph):
        if is_first_slide == True:
            slide = prs.slides[0]
            is_first_slide = False
        else:
            empty_slide_layout = prs.slide_layouts[EMPTY_SLIDE]
            slide = prs.slides.add_slide(empty_slide_layout)

        self.AddLyric(slide, song_lyric_paragraph)
        self.AddTitle(slide, title)


    def CreatePPT(self, song_title_list, song_dict):
        log.debug ("PPTProduct:CreatePPT, title:{} content:{} ".format(song_title_list, len(song_dict)))
        prs = Presentation("template.pptx")

        is_first_slide = True

        for title in song_title_list:
            song_lyric_paragraph_list = song_dict.get(title)
            if song_lyric_paragraph_list == None:
                continue

            for song_lyric_paragraph in song_lyric_paragraph_list:
                slide = ""


                if len(song_lyric_paragraph) > LANGEST_LINE_IN_PPT:
                    pages = len(song_lyric_paragraph) / LANGEST_LINE_IN_PPT
                    for index in range(0,pages):
                        if is_first_slide == True:
                            slide = prs.slides[0]
                            is_first_slide = False
                        else:
                            empty_slide_layout = prs.slide_layouts[EMPTY_SLIDE]
                            slide = prs.slides.add_slide(empty_slide_layout)

                            self.AddLyric(slide, song_lyric_paragraph[
                                index * LANGEST_LINE_IN_PPT: (index + 1) * LANGEST_LINE_IN_PPT])

                        self.AddTitle(slide, title)

                    empty_slide_layout = prs.slide_layouts[EMPTY_SLIDE]
                    slide = prs.slides.add_slide(empty_slide_layout)

                    self.AddLyric(slide, song_lyric_paragraph[ pages * LANGEST_LINE_IN_PPT: ])
                    self.AddTitle(slide, title)
                else:

                    if is_first_slide == True:
                        slide = prs.slides[0]
                        is_first_slide = False
                    else:
                        empty_slide_layout = prs.slide_layouts[EMPTY_SLIDE]
                        slide = prs.slides.add_slide(empty_slide_layout)

                    self.AddLyric(slide, song_lyric_paragraph)
                    self.AddTitle(slide, title)

        prs.save(today + '.pptx')
        log.debug ("PPTProduct:CreatePPT: {}.pptx is save".format(today))

    def TestCreatePPT(self):
        log.debug ("PPTProduct:TestCreatePPT")

        prs = Presentation("template.pptx")
        title_slide_layout = prs.slide_layouts[DEFAULT_SLIDE]
        slide = prs.slides.add_slide(title_slide_layout)

        title = slide.shapes.title
        shapes = slide.shapes
        subtitle = slide.placeholders[1]

        title.text = "Hello, World!\n test\n test\n"
        subtitle.text = "python-ppt was here!"

        # left = top = width = height = Cm(8.0)
        # slide.shapes.add_shape(
        #     MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        # )

        prs.save('test.pptx')
        log.debug ("PPTProduct:TestCreatePPT: test.ppt is save")

    def analyze_ppt(self, input, output):
        """ Take the input file and analyze the structure.
        The output file contains marked up information to make it easier
        for generating future powerpoint templates.
        """
        prs = Presentation(input)
        # Each powerpoint file has multiple layouts
        # Loop through them all and  see where the various elements are
        for index, _ in enumerate(prs.slide_layouts):
            slide = prs.slides.add_slide(prs.slide_layouts[index])
            # Not every slide has to have a title
            try:
                title = slide.shapes.title
                title.text = 'Title for Layout {}'.format(index)
            except AttributeError:
                print("No Title for Layout {}".format(index))
            # Go through all the placeholders and identify them by index and type
            for shape in slide.placeholders:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    # Do not overwrite the title which is just a special placeholder
                    try:
                        if 'Title' not in shape.text:
                            shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
                    except AttributeError:
                        print("{} has no text attribute".format(phf.type))
                    print('{} {}'.format(phf.idx, shape.name))
        prs.save(output)

    def Ping(self):
        try:
            ret_error="test ping"
            #ret_error = self.client.Ping("python")
            log.debug('Ping(), return:{}'.format(ret_error))
        except Exception:
            log.debug(('Get Exception Ping()%s'))

    def PingThread(self):
        while 1:
            log.debug('"test from py thread"')
            time.sleep(2)


    def Start(self):
        # Create two threads as follows
        try:
            thread.start_new_thread( self.PingThread, () )
        except:
            log.debug ("Error: unable to start thread")



def ReadSongTitleList(file_name):
    log.debug("ReadSongTitleList from file:{}".format(file_name.encode(ENCODING)))
    song_title_list = []
    f = open(file_name, "r")
    for line in f:
        if len(line[:-1]):
            song_tile = line.decode(ENCODING)[:-1]
            song_title_list.append(song_tile)
            log.info(u"读取输入歌曲目录:\t{}".format(song_tile))

    log.debug("ReadSongTitleList from file:{}, with {} song titles"
             .format(file_name.encode(ENCODING), len(song_title_list)))
    f.close()
    return song_title_list



if __name__ == "__main__":
    try:
        #thrift_client.Start()

        # song_title_list = [
        #     u'阿爸阿爸父',
        #     u'教会啊你要兴起',
        #     u'惟有主的话永长存'
        # ]

        log.info(u"开始读取:主日赞美诗歌名.txt...\n")
        song_title_list = ReadSongTitleList(u"主日赞美诗歌名.txt")
        log.info(u"读取:主日赞美诗歌名.txt 完成...\n")
        log.debug(song_title_list)

        ppt_product = PPTProduct()
        doc_product = DocProduct()

        log.info(u"开始读取:全部诗歌歌词.docx...\n")
        doc_product.ReadDoc(u"全部诗歌歌词.docx")
        log.info(u"读取:全部诗歌歌词.docx 完成...\n")

        log.info(u"开始查找:对应诗歌歌词...\n")
        (find_all, song_lyric_dict) = doc_product.GetSongList(song_title_list)
        log.info(u"查找对应诗歌歌词 完成...\n")

        if len(song_lyric_dict) > 0:
            product_ppt = True
            if ( find_all == False):
                log.info(u"未找到所有的歌词,建议将未找到歌词添加进 \"全部诗歌歌词.docx\" 再重新生成。")
                product_ppt == query_yes_no(u"是否继续生成歌词PPT 和 word?")

            if product_ppt == True:
                if find_all == False:
                    log.info(u"未找到所有歌词，开始生成对应PPT Word...\n")
                else:
                    log.info(u"找到所有歌词，开始生成对应PPT Word...\n")
                ppt_product.CreatePPT(song_title_list, song_lyric_dict)
                doc_product.WriteDoc(song_title_list, song_lyric_dict)
                log.info(u"生成对应PPT Word: {}.pptx {}.docx...\n".format(today, today))
        else:
            log.warning(u"找不到所有的歌词")
    except Exception:
        log.error("Got exception on TestThriftServer:%s", traceback.format_exc() )

    print u"输入任意键退出"
    raw_input()

